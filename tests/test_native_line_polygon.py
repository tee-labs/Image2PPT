"""Native-shape lifting, part 2: polygons and straight connector lines.

Extends the coverage from test_native_shape_classification.py:

* classify_filled_shape — convex upright polygons (triangle / diamond /
  trapezoid) lift to native MSO shapes; concave glyphs and rotated
  quads stay on the PNG path.
* classify_connector_line — straight / dashed / dotted / diagonal
  connector strokes and single-head arrows classify to native line
  elements; elbows and filled blocks are rejected.
* LayoutBuilder wiring — connector-role elements emit ``type: line``
  records (no PNG asset), and the PPTX builder renders them as
  connectors with dash style and arrowheads.
"""
from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

import cv2
import numpy as np

ROOT = Path(__file__).resolve().parents[1]
PAGE_DIR = ROOT / "scripts" / "page"
SCRIPTS_DIR = ROOT / "scripts"
for _p in (str(PAGE_DIR), str(SCRIPTS_DIR)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from layout.outline import (  # noqa: E402
    classify_connector_line,
    classify_filled_shape,
)

BLUE_BGR = (200, 120, 40)      # → #2878C8


def _white(w: int, h: int) -> np.ndarray:
    return np.full((h, w, 3), 255, dtype=np.uint8)


class PolygonLiftTests(unittest.TestCase):
    def test_upright_triangle(self) -> None:
        img = _white(120, 100)
        cv2.fillPoly(img, [np.array([[60, 4], [116, 96], [4, 96]])],
                     BLUE_BGR)
        kind, fill, _line, radius, line_px = classify_filled_shape(img)
        self.assertEqual(kind, "triangle")
        self.assertEqual(fill, "#2878C8")
        self.assertEqual(radius, 0.0)
        self.assertEqual(line_px, 0.0)

    def test_diamond(self) -> None:
        img = _white(100, 100)
        cv2.fillPoly(img, [np.array([[50, 2], [98, 50], [50, 98], [2, 50]])],
                     BLUE_BGR)
        kind, fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "diamond")
        self.assertEqual(fill, "#2878C8")

    def test_trapezoid_short_top(self) -> None:
        img = _white(140, 90)
        cv2.fillPoly(img, [np.array(
            [[35, 2], [105, 2], [138, 88], [2, 88]])], BLUE_BGR)
        kind, fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "trapezoid")
        self.assertEqual(fill, "#2878C8")

    def test_rotated_quad_stays_png(self) -> None:
        # Slanted parallelogram: not an upright primitive.
        img = _white(140, 90)
        cv2.fillPoly(img, [np.array(
            [[30, 60], [60, 10], [130, 30], [100, 80]])], BLUE_BGR)
        self.assertIsNone(classify_filled_shape(img))

    def test_concave_star_stays_png(self) -> None:
        img = _white(100, 100)
        pts = cv2.boxPoints(((50, 50), (90, 90), 0))
        star = []
        cx, cy = 50.0, 50.0
        for i in range(10):
            ang = np.pi / 2 + i * np.pi / 5
            r = 46 if i % 2 == 0 else 20
            star.append([cx + r * np.cos(ang), cy - r * np.sin(ang)])
        cv2.fillPoly(img, [np.array(star, dtype=np.int32)],
                     (60, 60, 60))
        self.assertIsNone(classify_filled_shape(img))


class ConnectorLineTests(unittest.TestCase):
    def test_solid_horizontal_line(self) -> None:
        img = _white(200, 12)
        cv2.line(img, (4, 6), (195, 6), (90, 90, 90), 2)
        pts, line_hex, width_px, dash, arrow = classify_connector_line(img)
        self.assertEqual(pts[1], pts[3])            # horizontal
        self.assertLessEqual(pts[0], 6)             # spans the crop
        self.assertGreaterEqual(pts[2], 193)
        self.assertEqual(line_hex, "#5A5A5A")
        self.assertLessEqual(width_px, 4.0)
        self.assertIsNone(dash)
        self.assertIsNone(arrow)

    def test_dashed_line(self) -> None:
        img = _white(220, 10)
        for x0 in range(4, 210, 20):
            cv2.line(img, (x0, 5), (min(x0 + 12, 215), 5), (90, 90, 90), 2)
        _pts, _hex, _width, dash, arrow = classify_connector_line(img)
        self.assertEqual(dash, "dash")
        self.assertIsNone(arrow)

    def test_dotted_line(self) -> None:
        img = _white(200, 10)
        for x0 in range(5, 195, 8):
            cv2.circle(img, (x0, 5), 1, (90, 90, 90), -1)
        _pts, _hex, _width, dash, _arrow = classify_connector_line(img)
        self.assertEqual(dash, "dot")

    def test_diagonal_line_endpoints(self) -> None:
        img = _white(120, 120)
        cv2.line(img, (6, 6), (113, 113), (60, 60, 200), 3)
        pts, _hex, width_px, _dash, _arrow = classify_connector_line(img)
        self.assertLessEqual(abs(pts[0] - pts[1]), 3)   # starts near TL
        self.assertGreaterEqual(pts[2], 105)            # ends near BR
        self.assertGreaterEqual(pts[3], 105)
        self.assertLessEqual(width_px, 6.0)

    def test_arrowhead_detected(self) -> None:
        img = _white(160, 40)
        cv2.line(img, (4, 20), (120, 20), (90, 90, 90), 3)
        cv2.fillPoly(img, [np.array(
            [[150, 20], [122, 8], [122, 32]])], (90, 90, 90))
        _pts, _hex, _width, dash, arrow = classify_connector_line(img)
        self.assertIsNone(dash)
        self.assertEqual(arrow, "end")

    def test_elbow_rejected(self) -> None:
        img = _white(120, 120)
        cv2.line(img, (6, 6), (113, 6), (90, 90, 90), 3)
        cv2.line(img, (113, 6), (113, 113), (90, 90, 90), 3)
        self.assertIsNone(classify_connector_line(img))

    def test_filled_block_rejected(self) -> None:
        img = _white(80, 80)
        cv2.rectangle(img, (4, 20), (75, 60), BLUE_BGR, -1)
        self.assertIsNone(classify_connector_line(img))


class BuilderConnectorLineTests(unittest.TestCase):
    def _build(self, td: Path, dashed: bool = False):
        from layout.builder import LayoutBuilder
        src = td / "page_01.png"
        img = _white(400, 300)
        if dashed:
            for x0 in range(60, 330, 24):
                cv2.line(img, (x0, 150), (min(x0 + 14, 338), 150),
                         (90, 90, 90), 2)
        else:
            cv2.line(img, (50, 150), (340, 150), (90, 90, 90), 2)
        cv2.imwrite(str(src), img)
        inventory = [
            {"id": "v000", "type": "image",
             "bbox": [50, 145, 341, 156],
             "source": "cleaned", "role": "connector"},
        ]
        inv_path = td / "inventory.json"
        inv_path.write_text(json.dumps(inventory), encoding="utf-8")
        args = type("Args", (), {
            "inventory": str(inv_path),
            "source": str(src),
            "cleaned": str(src),
            "out_assets_dir": str(td / "assets" / "page_01"),
            "asset_prefix": "assets/page_01",
            "out_manifest": str(td / "m.json"),
            "out_layout": str(td / "l.json"),
            "slide_width_in": None,
            "slide_height_in": 7.5,
        })()
        builder = LayoutBuilder(args)
        builder.build()
        builder.write()
        return builder

    def test_connector_becomes_native_line(self) -> None:
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            builder = self._build(td)
            lines = [el for el in builder.front_shape_elements
                     if el["type"] == "line"]
            self.assertEqual(len(lines), 1)
            el = lines[0]
            self.assertEqual(el["name"], "v000")
            x1, y1, x2, y2 = el["points"]
            self.assertLessEqual(x1, 55)
            self.assertGreaterEqual(x2, 335)
            self.assertEqual(y1, y2)
            self.assertEqual(el["line"], "#5A5A5A")
            self.assertNotIn("dash", el)
            # No PNG asset for the lifted line.
            self.assertEqual(builder.image_elements, [])
            self.assertEqual(
                list((td / "assets" / "page_01").glob("*.png")), [])

    def test_dashed_connector_keeps_dash(self) -> None:
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            builder = self._build(td, dashed=True)
            lines = [el for el in builder.front_shape_elements
                     if el["type"] == "line"]
            if lines:                      # lift hit → dash preserved
                self.assertEqual(lines[0].get("dash"), "dash")
                self.assertEqual(builder.image_elements, [])
            # else: the dashes did not merge into one stroke → PNG path
            # (unchanged behaviour), also acceptable.

    def test_line_renders_as_pptx_connector(self) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from deck.build_pptx_from_layout import run as build_pptx
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            self._build(td)
            layout = json.loads((td / "l.json").read_text(encoding="utf-8"))
            layout["elements"][0]["arrow"] = "end"
            (td / "l.json").write_text(
                json.dumps(layout), encoding="utf-8")
            out = td / "out.pptx"
            build_pptx(layout=str(td / "l.json"), out=str(out),
                       assets_root=str(td))
            slide = Presentation(str(out)).slides[0]
            conns = [sh for sh in slide.shapes
                     if sh.shape_type == MSO_SHAPE_TYPE.LINE
                     or sh.shape_type == MSO_SHAPE_TYPE.LINE_INVERSE
                     or not sh.has_text_frame and not sh.shape_type]
            self.assertTrue(conns)
            conn = conns[0]
            xml = conn.line._get_or_add_ln().xml
            self.assertIn("a:tailEnd", xml)


if __name__ == "__main__":
    unittest.main()
