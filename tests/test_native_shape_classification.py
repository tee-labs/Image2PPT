"""Native shape classification: circles / rects become editable shapes.

Regression coverage for the issue "原生形状没有做好": filled circles and
rectangles used to be flattened into PNG crops / backgrounds. These tests
lock in:

* classify_filled_shape — solid primitives (container/internal roles)
  classify to oval / round_rect / rect; gradients, photos and shapes
  with baked-in content stay on the PNG path.
* classify_outline_ring — near-square outline rings classify to oval or
  round_rect instead of being unconditionally excluded.
* LayoutBuilder wiring — native shape elements are emitted, z-ordered
  after images (front shapes) / before images (card frames), and the
  PPTX builder renders them as real AutoShapes.
"""
from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

import cv2
import numpy as np

ROOT = Path(__file__).resolve().parents[1]
PAGE_DIR = ROOT / "scripts" / "page"
SCRIPTS_DIR = ROOT / "scripts"
for _p in (str(PAGE_DIR), str(SCRIPTS_DIR)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from layout.outline import (  # noqa: E402
    classify_filled_shape,
    classify_outline_ring,
)

BLUE_BGR = (200, 120, 40)      # → #2878C8
RED_BGR = (60, 60, 180)        # → #B43C3C


def _white(size: int) -> np.ndarray:
    return np.full((size, size, 3), 255, dtype=np.uint8)


def _solid_circle(size: int = 96, color=BLUE_BGR) -> np.ndarray:
    img = _white(size)
    cv2.circle(img, (size // 2, size // 2), size // 2 - 2, color, -1)
    return img


def _rounded_square(size: int = 96, radius: int = 32,
                    color=BLUE_BGR) -> np.ndarray:
    img = _white(size)
    r = radius
    cv2.rectangle(img, (r, 0), (size - r, size), color, -1)
    cv2.rectangle(img, (0, r), (size, size - r), color, -1)
    for cx, cy in ((r, r), (size - r, r), (r, size - r), (size - r, size - r)):
        cv2.circle(img, (cx, cy), r, color, -1)
    return img


class ClassifyFilledShapeTests(unittest.TestCase):
    def test_solid_circle_is_oval(self) -> None:
        kind, fill, line, radius, line_px = classify_filled_shape(
            _solid_circle())
        self.assertEqual(kind, "oval")
        self.assertEqual(radius, 0.0)
        self.assertEqual(line_px, 0.0)
        self.assertEqual(fill, "#2878C8")
        self.assertTrue(line.startswith("#"))

    def test_solid_square_is_rect(self) -> None:
        img = np.full((80, 80, 3), 255, dtype=np.uint8)
        cv2.rectangle(img, (0, 0), (79, 79), RED_BGR, -1)
        kind, fill, _line, radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "rect")
        self.assertEqual(radius, 0.0)
        self.assertEqual(fill, "#B43C3C")

    def test_rounded_square_is_round_rect(self) -> None:
        kind, _fill, _line, radius, _line_px = classify_filled_shape(
            _rounded_square(radius=32))
        self.assertEqual(kind, "round_rect")
        self.assertGreaterEqual(radius, 0.08)
        self.assertLessEqual(radius, 0.5)

    def test_thin_ring_is_transparent_oval(self) -> None:
        # The issue's uploaded asset: a thin light-gray circle outline
        # (~86x88) on a white background.
        img = _white(96)
        cv2.circle(img, (48, 48), 44, (224, 224, 224), 4)
        kind, fill, line, radius, line_px = classify_filled_shape(img)
        self.assertEqual(kind, "oval")
        self.assertIsNone(fill)          # transparent interior
        self.assertEqual(line, "#E0E0E0")
        self.assertEqual(radius, 0.0)
        self.assertGreaterEqual(line_px, 2.0)
        self.assertLessEqual(line_px, 12.0)

    def test_donut_ring_is_transparent_oval(self) -> None:
        img = _white(120)
        cv2.circle(img, (60, 60), 56, BLUE_BGR, -1)
        cv2.circle(img, (60, 60), 28, (255, 255, 255), -1)
        kind, fill, line, _radius, line_px = classify_filled_shape(img)
        self.assertEqual(kind, "oval")
        self.assertIsNone(fill)
        self.assertEqual(line, "#2878C8")
        self.assertGreaterEqual(line_px, 10.0)

    def test_gradient_ring_rejected(self) -> None:
        img = _white(96)
        for t in range(88):
            col = (int(40 + 2 * t), int(120 + 0.5 * t), 200 - t)
            cv2.circle(img, (48, 48), 44, col, 1)
        self.assertIsNone(classify_filled_shape(img))

    def test_broken_ring_rejected(self) -> None:
        img = _white(96)
        cv2.circle(img, (48, 48), 44, BLUE_BGR, 4)
        cv2.rectangle(img, (40, 0), (56, 20), (255, 255, 255), -1)  # notch
        self.assertIsNone(classify_filled_shape(img))

    def test_offcentre_hole_rejected(self) -> None:
        # Filled disc with an off-centre punched hole: not a clean ring,
        # and its interior isn't uniform → stays a PNG.
        img = _white(120)
        cv2.circle(img, (60, 60), 56, BLUE_BGR, -1)
        cv2.circle(img, (40, 40), 20, (255, 255, 255), -1)
        self.assertIsNone(classify_filled_shape(img))

    def test_gradient_rect_rejected(self) -> None:
        img = np.zeros((60, 160, 3), dtype=np.uint8)
        top = np.array([220, 90, 30], dtype=np.float32)
        bottom = np.array([30, 60, 200], dtype=np.float32)
        for y in range(60):
            t = y / 59.0
            img[y, :] = (top * (1 - t) + bottom * t).astype(np.uint8)
        self.assertIsNone(classify_filled_shape(img))

    def test_photo_noise_rejected(self) -> None:
        rng = np.random.default_rng(7)
        img = rng.integers(0, 255, size=(80, 80, 3), dtype=np.uint8)
        self.assertIsNone(classify_filled_shape(img))

    def test_glyph_inside_circle_rejected(self) -> None:
        img = _solid_circle()
        cv2.circle(img, (48, 48), 12, (255, 255, 255), -1)
        self.assertIsNone(classify_filled_shape(img))

    def test_circle_on_tinted_card(self) -> None:
        img = np.full((96, 96, 3), (235, 245, 250), dtype=np.uint8)
        cv2.circle(img, (48, 48), 46, BLUE_BGR, -1)
        kind, fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "oval")
        self.assertEqual(fill, "#2878C8")

    def test_full_bleed_square_is_rect(self) -> None:
        img = np.full((70, 110, 3), RED_BGR, dtype=np.uint8)
        kind, _fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "rect")

    def test_tiny_crop_rejected(self) -> None:
        self.assertIsNone(classify_filled_shape(
            np.full((8, 8, 3), 120, dtype=np.uint8)))

    def test_dashed_pale_placeholder_rejected(self) -> None:
        # Ghost-pale blocks (dashed photo placeholders, faint tints)
        # keep the PNG path: a flat native fill would lose the borders.
        img = np.full((120, 200, 3), (244, 244, 244), dtype=np.uint8)
        for x in range(10, 190, 12):
            cv2.line(img, (x, 10), (x + 6, 10), (180, 180, 180), 1)
        self.assertIsNone(classify_filled_shape(img))

    def test_pale_bordered_card_becomes_native_rect(self) -> None:
        # The classic deck card: ghost-pale uniform fill + a visible
        # 1 px border. With a crisp border present it belongs in the
        # native-shape path (fill + line both preserved); only truly
        # borderless faint tints keep the PNG path.
        img = np.full((200, 300, 3), (238, 245, 250), dtype=np.uint8)
        cv2.rectangle(img, (0, 0), (299, 199), (200, 200, 200), 1)
        kind, fill, line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "rect")
        self.assertEqual(fill, "#FAF5EE")
        self.assertEqual(line, "#C8C8C8")

    def test_pale_borderless_tint_rejected(self) -> None:
        # No border stroke anywhere: a flat native fill would render an
        # invisible box over a soft-edged tint, so keep the PNG path.
        img = np.full((200, 300, 3), (238, 245, 250), dtype=np.uint8)
        self.assertIsNone(classify_filled_shape(img))

    def test_user_uploaded_ring_asset(self) -> None:
        """The exact 86x88 RGBA asset from the issue, composited on
        white the way convert.py ingest now does."""
        path = Path(__file__).parent / "fixtures" / "issue5_ring.png"
        if not path.exists():
            self.skipTest("fixture not present")
        rgba = cv2.imread(str(path), cv2.IMREAD_UNCHANGED)
        alpha = rgba[:, :, 3:4].astype(np.float32) / 255.0
        comp = (rgba[:, :, :3].astype(np.float32) * alpha
                + 255.0 * (1.0 - alpha)).astype(np.uint8)
        kind, fill, line, _radius, line_px = classify_filled_shape(comp)
        self.assertEqual(kind, "oval")
        self.assertIsNone(fill)
        self.assertGreater(line_px, 1.0)


def _ring_mask(size_h: int, size_w: int, center: tuple[int, int],
               radius: int, thickness: int = 5) -> np.ndarray:
    mask = np.zeros((size_h, size_w), dtype=np.uint8)
    cv2.circle(mask, center, radius, 255, thickness)
    return cv2.dilate(mask, np.ones((2, 2), np.uint8), iterations=1)


class ClassifyOutlineRingTests(unittest.TestCase):
    def test_circle_ring_is_oval(self) -> None:
        source = _white(200)
        cv2.circle(source, (100, 100), 70, (180, 100, 40), 4)
        mask = _ring_mask(200, 200, (100, 100), 70)
        with tempfile.TemporaryDirectory() as td:
            mask_path = Path(td) / "v.mask.png"
            cv2.imwrite(str(mask_path), mask)
            result = classify_outline_ring(
                source, (30, 30, 170, 170), str(mask_path))
        self.assertIsNotNone(result)
        kind, radius = result
        self.assertEqual(kind, "oval")
        self.assertEqual(radius, 0.0)

    def test_square_ring_is_rect(self) -> None:
        # A sharp-cornered frame lifts as a plain rect (phase 3: the old
        # 0.08 minimum radius rounded corners the source does not have).
        source = _white(200)
        cv2.rectangle(source, (40, 40), (160, 160), (180, 100, 40), 4)
        mask = np.zeros((200, 200), dtype=np.uint8)
        cv2.rectangle(mask, (40, 40), (160, 160), 255, 5)
        mask = cv2.dilate(mask, np.ones((2, 2), np.uint8), iterations=1)
        with tempfile.TemporaryDirectory() as td:
            mask_path = Path(td) / "v.mask.png"
            cv2.imwrite(str(mask_path), mask)
            result = classify_outline_ring(
                source, (40, 40, 160, 160), str(mask_path))
        self.assertIsNotNone(result)
        self.assertEqual(result, ("rect", 0.0))

    def test_concentric_rings_rejected(self) -> None:
        source = _white(200)
        cv2.circle(source, (100, 100), 70, (180, 100, 40), 4)
        cv2.circle(source, (100, 100), 36, (180, 100, 40), 4)
        mask = _ring_mask(200, 200, (100, 100), 70)
        cv2.circle(mask, (100, 100), 36, 255, 5)
        with tempfile.TemporaryDirectory() as td:
            mask_path = Path(td) / "v.mask.png"
            cv2.imwrite(str(mask_path), mask)
            result = classify_outline_ring(
                source, (30, 30, 170, 170), str(mask_path))
        self.assertIsNone(result)

    def test_content_polluted_mask_rejected(self) -> None:
        source = _white(200)
        cv2.circle(source, (100, 100), 70, (180, 100, 40), 4)
        rng = np.random.default_rng(3)
        source[75:125, 75:125] = rng.integers(
            0, 200, size=(50, 50, 3), dtype=np.uint8)
        result = classify_outline_ring(source, (30, 30, 170, 170), None)
        self.assertIsNone(result)

    def test_small_ring_rejected(self) -> None:
        source = _white(60)
        cv2.circle(source, (30, 30), 16, (180, 100, 40), 3)
        self.assertIsNone(classify_outline_ring(source, (14, 14, 46, 46)))


def _write_slide(path: Path) -> None:
    img = np.full((540, 960, 3), 255, dtype=np.uint8)
    cv2.circle(img, (128, 128), 48, BLUE_BGR, -1)          # container circle
    cv2.rectangle(img, (220, 300), (319, 399), RED_BGR, -1)  # internal square
    # internal gradient block → must stay a PNG
    top = np.array([220, 90, 30], dtype=np.float32)
    bottom = np.array([30, 60, 200], dtype=np.float32)
    grad = np.zeros((100, 100, 3), dtype=np.uint8)
    for y in range(100):
        t = y / 99.0
        grad[y, :] = (top * (1 - t) + bottom * t).astype(np.uint8)
    img[300:400, 600:700] = grad
    cv2.circle(img, (480, 140), 50, (180, 100, 40), 4)     # outline ring
    cv2.imwrite(str(path), img)


def _write_fg_mask(path: Path, img: np.ndarray,
                   bbox: tuple[int, int, int, int]) -> None:
    """Real-run style mask: foreground pixels opaque inside the bbox."""
    mask = np.zeros(img.shape[:2], dtype=np.uint8)
    x1, y1, x2, y2 = bbox
    crop = img[y1:y2, x1:x2]
    fg = np.abs(crop.astype(int) - 255).max(axis=2) > 14
    mask[y1:y2, x1:x2] = (fg * 255).astype(np.uint8)
    cv2.imwrite(str(path), mask)


def _write_ring_mask(path: Path) -> None:
    cv2.imwrite(str(path), _ring_mask(540, 960, (480, 140), 50))


class BuilderNativeShapeTests(unittest.TestCase):
    def _build(self, td: Path):
        from layout.builder import LayoutBuilder
        src = td / "page_01.png"
        _write_slide(src)
        masks = td / "masks"
        masks.mkdir(exist_ok=True)
        slide = cv2.imread(str(src))
        # Real runs write per-element masks; the builder must classify
        # the crop BEFORE the masked path can swallow it as a PNG.
        for el_id, bbox in (("v000", (80, 80, 176, 176)),
                            ("v001", (220, 300, 320, 400)),
                            ("v002", (600, 300, 700, 400))):
            _write_fg_mask(masks / f"{el_id}.mask.png", slide, bbox)
        mask_path = masks / "v003.mask.png"
        _write_ring_mask(mask_path)
        inventory = [
            {"id": "v000", "type": "image",
             "bbox": [80, 80, 176, 176],
             "source": "cleaned", "role": "container",
             "mask_path": str(masks / "v000.mask.png")},
            {"id": "v001", "type": "image",
             "bbox": [220, 300, 320, 400],
             "source": "source", "role": "internal",
             "mask_path": str(masks / "v001.mask.png")},
            {"id": "v002", "type": "image",
             "bbox": [600, 300, 700, 400],
             "source": "source", "role": "internal",
             "mask_path": str(masks / "v002.mask.png")},
            {"id": "v003", "type": "image",
             "bbox": [430, 90, 530, 190],
             "source": "cleaned", "role": "outline",
             "mask_path": str(mask_path)},
        ]
        inv_path = td / "inventory.json"
        inv_path.write_text(json.dumps(inventory), encoding="utf-8")
        args = type("Args", (), {
            "inventory": str(inv_path),
            "source": str(src),
            "cleaned": str(src),
            "out_assets_dir": str(td / "assets" / "page_01"),
            "asset_prefix": "assets/page_01",
            "out_manifest": str(td / "m.json"),
            "out_layout": str(td / "l.json"),
            "slide_width_in": None,
            "slide_height_in": 7.5,
        })()
        builder = LayoutBuilder(args)
        builder.build()
        builder.write()
        return builder

    def test_solid_shapes_become_front_native_shapes(self) -> None:
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            builder = self._build(td)
            fronts = {el["name"]: el for el in builder.front_shape_elements}
            self.assertEqual(set(fronts), {"v000", "v001"})
            self.assertEqual(fronts["v000"]["shape"], "oval")
            self.assertEqual(fronts["v000"]["box"], [80, 80, 96, 96])
            self.assertEqual(fronts["v001"]["shape"], "rect")
            self.assertEqual(fronts["v001"]["box"], [220, 300, 100, 100])
            # No PNG assets / image elements for the lifted shapes.
            names = {el["name"] for el in builder.image_elements}
            self.assertEqual(names, {"v002"})
            assets = {p.name
                      for p in (td / "assets" / "page_01").glob("*.png")}
            self.assertEqual(assets, {"v002.png"})

    def test_near_square_outline_ring_becomes_oval(self) -> None:
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            builder = self._build(td)
            self.assertEqual(len(builder.shape_elements), 1)
            shape = builder.shape_elements[0]
            self.assertEqual(shape["name"], "v003")
            self.assertEqual(shape["shape"], "oval")
            self.assertEqual(shape["box"], [430, 90, 100, 100])
            self.assertTrue(shape["line"])

    def test_z_order_shapes_images_front_shapes_text(self) -> None:
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            builder = self._build(td)
            layout = json.loads(
                (td / "l.json").read_text(encoding="utf-8"))
            kinds = [el["type"] for el in layout["elements"]]
            self.assertEqual(
                kinds,
                ["shape", "image", "shape", "shape"])
            self.assertEqual(layout["elements"][0]["name"], "v003")
            self.assertEqual(layout["elements"][1]["name"], "v002")
            self.assertEqual(layout["elements"][2]["name"], "v000")
            self.assertEqual(layout["elements"][3]["name"], "v001")

    def test_roleless_ring_asset_becomes_native_oval(self) -> None:
        """The issue's exact upload: a thin light-gray circle outline,
        inventoried with role=None (it is neither card nor connector).
        Even with a mask present it must lift to a transparent-fill
        native oval, not a flattened PNG."""
        from layout.builder import LayoutBuilder
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            src = td / "page_01.png"
            img = np.full((540, 960, 3), 255, dtype=np.uint8)
            cv2.circle(img, (200, 150), 42, (224, 224, 224), 4)
            cv2.imwrite(str(src), img)
            masks = td / "masks"
            masks.mkdir()
            _write_fg_mask(masks / "v000.mask.png", img, (154, 104, 246, 196))
            inventory = [
                {"id": "v000", "type": "image",
                 "bbox": [154, 104, 246, 196],
                 "source": "cleaned", "role": None,
                 "mask_path": str(masks / "v000.mask.png")},
            ]
            inv_path = td / "inventory.json"
            inv_path.write_text(json.dumps(inventory), encoding="utf-8")
            args = type("Args", (), {
                "inventory": str(inv_path),
                "source": str(src),
                "cleaned": str(src),
                "out_assets_dir": str(td / "assets" / "page_01"),
                "asset_prefix": "assets/page_01",
                "out_manifest": str(td / "m.json"),
                "out_layout": str(td / "l.json"),
                "slide_width_in": None,
                "slide_height_in": 7.5,
            })()
            builder = LayoutBuilder(args)
            builder.build()
            builder.write()
            self.assertEqual(len(builder.front_shape_elements), 1)
            shape = builder.front_shape_elements[0]
            self.assertEqual(shape["shape"], "oval")
            self.assertEqual(shape["fill"], "transparent")
            self.assertEqual(shape["line"], "#E0E0E0")
            self.assertEqual(shape["box"], [154, 104, 92, 92])
            # ring is ~4px @ 540-tall on a 7.5in slide → ~3pt stroke
            self.assertGreaterEqual(shape["line_width"], 2.0)
            self.assertEqual(builder.image_elements, [])

    def test_rgba_source_composited_on_white_at_ingest(self) -> None:
        """convert.copy_as_pages must composite transparent PNGs onto
        white — cv2's alpha-dropping imread would show black."""
        from convert import copy_as_pages
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            rgba = np.zeros((40, 40, 4), dtype=np.uint8)
            rgba[:, :, :3] = (40, 120, 200)
            rgba[10:30, 10:30, 3] = 255
            src = td / "in.png"
            cv2.imwrite(str(src), rgba)
            out_dir = copy_as_pages([src], td / "pages")
            page = cv2.imread(str(next(out_dir.glob("page_*.png"))))
            self.assertEqual(page.shape, (40, 40, 3))       # alpha dropped
            self.assertTrue((page[0, 0] >= 250).all())       # white bg
            self.assertAlmostEqual(int(page[20, 20][0]), 40, delta=2)

    def test_layout_builds_pptx_with_native_autoshapes(self) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from deck.build_pptx_from_layout import run as build_pptx
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            self._build(td)
            out = td / "out.pptx"
            build_pptx(layout=str(td / "l.json"), out=str(out),
                       assets_root=str(td))
            prs = Presentation(str(out))
            slide = prs.slides[0]
            autos = [sh for sh in slide.shapes
                     if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
            pics = [sh for sh in slide.shapes
                    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
            self.assertEqual(len(autos), 3)
            self.assertEqual(len(pics), 1)
            from pptx.enum.shapes import MSO_SHAPE
            kinds = sorted(sh.auto_shape_type for sh in autos)
            ovals = [k for k in kinds if k == MSO_SHAPE.OVAL]
            rects = [k for k in kinds if k == MSO_SHAPE.RECTANGLE]
            self.assertEqual(len(ovals), 2)
            self.assertEqual(len(rects), 1)

    def test_badge_with_erased_text_becomes_oval(self) -> None:
        """The real badge flow: glyph erased in text_only, so the badge
        classifies to a native oval and the OCR text renders on top."""
        from layout.builder import LayoutBuilder
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            src = td / "page_01.png"
            img = np.full((540, 960, 3), 255, dtype=np.uint8)
            cv2.circle(img, (100, 100), 48, BLUE_BGR, -1)
            cv2.putText(img, "3", (88, 116), cv2.FONT_HERSHEY_SIMPLEX,
                        1.3, (255, 255, 255), 3, cv2.LINE_AA)
            cv2.imwrite(str(src), img)
            # Real work-dir layout: cleaned/<stem>.png + <stem>.text_only.png
            clean_dir = td / "clean"
            clean_dir.mkdir()
            clean = np.full((540, 960, 3), 255, dtype=np.uint8)
            cv2.circle(clean, (100, 100), 48, BLUE_BGR, -1)
            cv2.imwrite(str(clean_dir / "page_01.png"), clean)
            cv2.imwrite(str(clean_dir / "page_01.text_only.png"), clean)
            inventory = [
                {"id": "t000", "type": "text", "text": "3",
                 "bbox": [88, 84, 116, 118], "confidence": 0.95},
                {"id": "v000", "type": "image",
                 "bbox": [52, 52, 148, 148],
                 "source": "source", "role": "internal"},
            ]
            inv_path = td / "inventory.json"
            inv_path.write_text(json.dumps(inventory), encoding="utf-8")
            args = type("Args", (), {
                "inventory": str(inv_path),
                "source": str(src),
                "cleaned": str(clean_dir / "page_01.png"),
                "out_assets_dir": str(td / "assets" / "page_01"),
                "asset_prefix": "assets/page_01",
                "out_manifest": str(td / "m.json"),
                "out_layout": str(td / "l.json"),
                "slide_width_in": None,
                "slide_height_in": 7.5,
            })()
            builder = LayoutBuilder(args)
            builder.build()
            builder.write()
            self.assertEqual(len(builder.front_shape_elements), 1)
            shape = builder.front_shape_elements[0]
            self.assertEqual(shape["shape"], "oval")
            self.assertEqual(shape["box"], [52, 52, 96, 96])
            self.assertEqual(builder.image_elements, [])

    def test_variety_of_primitives_map_to_matching_shapes(self) -> None:
        from layout.builder import LayoutBuilder
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            src = td / "page_01.png"
            img = np.full((540, 960, 3), 255, dtype=np.uint8)
            cv2.ellipse(img, (60, 500), (55, 30), 0, 0, 360,
                        (150, 80, 200), -1)              # wide ellipse
            cv2.rectangle(img, (300, 420), (699, 499),
                          (90, 90, 90), -1)              # wide banner
            r0 = 14
            chip = np.full((540, 960, 3), 255, dtype=np.uint8)
            cv2.rectangle(chip, (760 + r0, 430), (940 - r0, 500),
                          (40, 180, 90), -1)
            cv2.rectangle(chip, (760, 430 + r0), (940, 500 - r0),
                          (40, 180, 90), -1)
            for cx, cy in ((760 + r0, 430 + r0), (940 - r0, 430 + r0),
                           (760 + r0, 500 - r0), (940 - r0, 500 - r0)):
                cv2.circle(chip, (cx, cy), r0, (40, 180, 90), -1)
            img[430:501, 760:941] = chip[430:501, 760:941]  # rounded chip
            cv2.imwrite(str(src), img)
            inventory = [
                {"id": "v000", "type": "image",
                 "bbox": [5, 470, 115, 530],
                 "source": "cleaned", "role": "container"},
                {"id": "v001", "type": "image",
                 "bbox": [300, 420, 700, 500],
                 "source": "cleaned", "role": "container"},
                {"id": "v002", "type": "image",
                 "bbox": [760, 430, 941, 501],
                 "source": "cleaned", "role": "container"},
            ]
            inv_path = td / "inventory.json"
            inv_path.write_text(json.dumps(inventory), encoding="utf-8")
            args = type("Args", (), {
                "inventory": str(inv_path),
                "source": str(src),
                "cleaned": str(src),
                "out_assets_dir": str(td / "assets" / "page_01"),
                "asset_prefix": "assets/page_01",
                "out_manifest": str(td / "m.json"),
                "out_layout": str(td / "l.json"),
                "slide_width_in": None,
                "slide_height_in": 7.5,
            })()
            builder = LayoutBuilder(args)
            builder.build()
            builder.write()
            by_name = {el["name"]: el["shape"]
                       for el in builder.front_shape_elements}
            self.assertEqual(by_name,
                             {"v000": "oval", "v001": "rect",
                              "v002": "round_rect"})


if __name__ == "__main__":
    unittest.main()
