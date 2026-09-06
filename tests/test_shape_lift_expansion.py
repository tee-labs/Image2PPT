"""Native shape lift expansion: capsules, ellipse rings, skew polygons,
rotated rects.

Optimization pass over native shape recognition ("原生形状识别优化"):

* round-rect radius normalised by min side (was sqrt(area)) — wide
  rounded banners/capsules no longer measure a fraction of their true
  corner radius; a capsule now reaches the full 0.5.
* classify_outline_ring handles any aspect: wide rect rings keep their
  round_rect, wide ellipse rings lift to oval (they used to be blind
  round_rects drawn over the ellipse).
* classify_filled_shape lifts wide thin ellipse ring crops (title
  halos) that both the circle-ring and quad-ring candidates rejected.
* _polygon_candidate lifts parallelograms and orientation-matched
  regular pentagons / hexagons.
* _rotated_rect_candidate lifts skewed solid rectangles with a
  rotation the layout carries into the PPTX autoshape.
"""
from __future__ import annotations

import json
import math
import sys
import tempfile
import unittest
from pathlib import Path

import cv2
import numpy as np

ROOT = Path(__file__).resolve().parents[1]
PAGE_DIR = ROOT / "scripts" / "page"
SCRIPTS_DIR = ROOT / "scripts"
for _p in (str(PAGE_DIR), str(SCRIPTS_DIR)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from layout.outline import (  # noqa: E402
    _rotated_rect_candidate,
    classify_filled_shape,
    classify_outline_ring,
)

BLUE_BGR = (200, 120, 40)      # → #2878C8
GRAY_BGR = (90, 90, 90)        # → #5A5A5A


def _white(w: int, h: int) -> np.ndarray:
    return np.full((h, w, 3), 255, dtype=np.uint8)


def _rounded_rect_img(w: int, h: int, r: int,
                      color=BLUE_BGR) -> np.ndarray:
    img = _white(w + 8, h + 8)
    cv2.rectangle(img, (4 + r, 4), (4 + w - r, 4 + h), color, -1)
    cv2.rectangle(img, (4, 4 + r), (4 + w, 4 + h - r), color, -1)
    for cx, cy in ((4 + r, 4 + r), (4 + w - r, 4 + r),
                   (4 + r, 4 + h - r), (4 + w - r, 4 + h - r)):
        cv2.circle(img, (cx, cy), r, color, -1)
    return img


class RadiusNormalizationTests(unittest.TestCase):
    def test_capsule_reaches_full_rounding(self) -> None:
        """A 3:1 capsule (fully semicircular ends) must emit radius 0.5,
        not the sqrt(area)-normalised 0.29 the old formula produced."""
        kind, _fill, _line, radius, _line_px = classify_filled_shape(
            _rounded_rect_img(300, 100, 50))
        self.assertEqual(kind, "round_rect")
        self.assertGreaterEqual(radius, 0.48)
        self.assertLessEqual(radius, 0.5)

    def test_wide_rounded_rect_radius_matches_min_side(self) -> None:
        """A 4:1 rounded banner with r = 0.25 * h must measure ~0.25
        (the MSO adjustment is r / min(w, h))."""
        kind, _fill, _line, radius, _line_px = classify_filled_shape(
            _rounded_rect_img(400, 100, 25))
        self.assertEqual(kind, "round_rect")
        self.assertAlmostEqual(radius, 0.25, delta=0.06)

    def test_near_square_rounded_square_unchanged(self) -> None:
        kind, _fill, _line, radius, _line_px = classify_filled_shape(
            _rounded_rect_img(96, 96, 32))
        self.assertEqual(kind, "round_rect")
        self.assertAlmostEqual(radius, 32 / 96, delta=0.05)


class WideOutlineRingTests(unittest.TestCase):
    def test_wide_rect_ring_outline_is_round_rect(self) -> None:
        source = _white(720, 400)
        cv2.rectangle(source, (60, 100), (659, 339), (205, 160, 90), 3)
        result = classify_outline_ring(source, (60, 100, 660, 340))
        self.assertIsNotNone(result)
        kind, radius = result
        self.assertEqual(kind, "round_rect")
        self.assertGreaterEqual(radius, 0.08)
        self.assertLessEqual(radius, 0.5)

    def test_wide_ellipse_ring_outline_is_oval(self) -> None:
        """A wide ellipse outline used to get a blind round_rect drawn
        over it (extra visible box); it must lift to a native oval."""
        source = _white(500, 300)
        cv2.ellipse(source, (250, 150), (200, 90), 0, 0, 360,
                    (180, 100, 40), 4)
        result = classify_outline_ring(source, (50, 60, 450, 240))
        self.assertIsNotNone(result)
        kind, radius = result
        self.assertEqual(kind, "oval")
        self.assertEqual(radius, 0.0)

    def test_skewed_ellipse_ring_outline_is_oval(self) -> None:
        source = _white(500, 340)
        cv2.ellipse(source, (250, 170), (190, 80), 30, 0, 360,
                    (180, 100, 40), 4)
        pts = cv2.findNonZero(
            (np.abs(source.astype(int) - 255).max(axis=2) > 14)
            .astype(np.uint8))
        x, y, w, h = cv2.boundingRect(pts)
        result = classify_outline_ring(
            source, (int(x), int(y), int(x + w), int(y + h)))
        self.assertIsNotNone(result)
        self.assertEqual(result[0], "oval")


class EllipseRingCropTests(unittest.TestCase):
    def test_wide_ellipse_ring_crop_lifts(self) -> None:
        """Wide thin ellipse ring crops (title halos) used to stay PNG
        crops: _ring_candidate needs a near-square annulus and
        _quad_ring_candidate needs straight sides."""
        img = _white(284, 128)
        cv2.ellipse(img, (142, 64), (140, 62), 0, 0, 360,
                    (180, 100, 40), 4)
        kind, fill, line, radius, line_px = classify_filled_shape(img)
        self.assertEqual(kind, "oval")
        self.assertIsNone(fill)          # transparent interior
        self.assertEqual(line, "#2864B4")
        self.assertEqual(radius, 0.0)
        self.assertGreaterEqual(line_px, 2.0)

    def test_solid_wide_ellipse_not_ring(self) -> None:
        img = _white(284, 128)
        cv2.ellipse(img, (142, 64), (140, 62), 0, 0, 360, BLUE_BGR, -1)
        kind, fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "oval")
        self.assertIsNotNone(fill)       # solid fill, not a ring

    def test_concentric_ellipse_rings_rejected(self) -> None:
        img = _white(284, 128)
        cv2.ellipse(img, (142, 64), (140, 62), 0, 0, 360,
                    (180, 100, 40), 4)
        cv2.ellipse(img, (142, 64), (60, 22), 0, 0, 360,
                    (180, 100, 40), 3)
        self.assertIsNone(classify_filled_shape(img))


class PolygonExpansionTests(unittest.TestCase):
    def test_parallelogram_lifts(self) -> None:
        img = _white(240, 140)
        pts = np.array([[50, 0], [230, 0], [200, 140], [20, 140]],
                       np.int32)
        cv2.fillPoly(img, [pts], GRAY_BGR)
        kind, _fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "parallelogram")

    def test_rectangle_not_parallelogram(self) -> None:
        img = _white(200, 100)
        cv2.rectangle(img, (0, 0), (199, 99), GRAY_BGR, -1)
        kind, _fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "rect")

    def test_regular_pentagon_lifts(self) -> None:
        # Canvas matches the pentagon bbox (1.902R x 1.809R) + 2px pad:
        # classify_filled_shape expects a tight bbox crop.
        r = 70
        img = _white(int(2 * r * math.sin(math.radians(72))) + 4,
                     int(r * (1 + math.sin(math.radians(54)))) + 4)
        cx, cy = img.shape[1] / 2, 2 + r
        pts = []
        for i in range(5):
            a = -math.pi / 2 + i * 2 * math.pi / 5
            pts.append((cx + r * math.cos(a), cy + r * math.sin(a)))
        cv2.fillPoly(img, [np.array(pts, np.int32)], GRAY_BGR)
        kind, _fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "pentagon")

    def test_rotated_pentagon_stays_png(self) -> None:
        """MSO REGULAR_PENTAGON draws point-up; a 30°-turned pentagon
        would render at the wrong angle, so it must stay a PNG."""
        r = 80
        img = _white(200, 200)
        pts = []
        for i in range(5):
            a = -math.pi / 2 + i * 2 * math.pi / 5 + math.radians(30)
            pts.append((100 + r * math.cos(a), 100 + r * math.sin(a)))
        cv2.fillPoly(img, [np.array(pts, np.int32)], GRAY_BGR)
        self.assertIsNone(classify_filled_shape(img))

    def test_flat_top_hexagon_lifts(self) -> None:
        # Canvas matches the hexagon bbox (2R x 1.732R) + 2px pad.
        r = 95
        img = _white(2 * r + 4, int(2 * r * math.sin(math.pi / 3)) + 4)
        cx, cy = img.shape[1] / 2, img.shape[0] / 2
        pts = []
        for i in range(6):
            a = i * math.pi / 3        # points left/right, flat top
            pts.append((cx + r * math.cos(a), cy + r * math.sin(a)))
        cv2.fillPoly(img, [np.array(pts, np.int32)], GRAY_BGR)
        kind, _fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "hexagon")


class RotatedRectTests(unittest.TestCase):
    def _rotated_rect_img(self, angle_deg: float, w: int = 200,
                          h: int = 70) -> np.ndarray:
        # Tight canvas around the rotated bbox: classify/rot candidates
        # expect bbox crops, not scenes with large margins.
        th = math.radians(angle_deg)
        c, s = math.cos(th), math.sin(th)
        bw = int(abs(w * c) + abs(h * s)) + 8
        bh = int(abs(w * s) + abs(h * c)) + 8
        img = _white(bw, bh)
        corners = np.array(
            [(-w / 2, -h / 2), (w / 2, -h / 2),
             (w / 2, h / 2), (-w / 2, h / 2)], np.float64)
        R = np.array([[c, -s], [s, c]])
        pts = (corners @ R.T) + np.array([bw / 2, bh / 2])
        cv2.fillPoly(img, [np.round(pts).astype(np.int32)], BLUE_BGR)
        return img

    def test_rotated_rect_lifts_with_rotation(self) -> None:
        img = self._rotated_rect_img(30)
        hit = _rotated_rect_candidate(img)
        self.assertIsNotNone(hit)
        x, y, w, h, rotation, fill, line = hit
        self.assertAlmostEqual(w / h, 200 / 70, delta=0.3)
        self.assertAlmostEqual(abs(rotation), 30.0, delta=3.0)
        self.assertEqual(fill, "#2878C8")

    def test_negative_rotation_sign_preserved(self) -> None:
        img = self._rotated_rect_img(-25)
        hit = _rotated_rect_candidate(img)
        self.assertIsNotNone(hit)
        rotation = hit[4]
        self.assertAlmostEqual(rotation, -25.0, delta=3.0)

    def test_upright_rect_not_captured(self) -> None:
        """Near-axis silhouettes belong to the upright bands; the
        rotated path must abstain (|angle| < 12° gate)."""
        img = _white(240, 80)
        cv2.rectangle(img, (0, 0), (239, 79), BLUE_BGR, -1)
        self.assertIsNone(_rotated_rect_candidate(img))

    def test_circle_not_captured(self) -> None:
        """A rasterized circle's min-area rect is a tilted square with
        IoU up to ~0.95; the sharp-corner gate must reject it."""
        img = _white(164, 164)
        cv2.circle(img, (82, 82), 80, BLUE_BGR, -1)
        self.assertIsNone(_rotated_rect_candidate(img))

    def test_upright_classifier_still_handles_upright_rect(self) -> None:
        img = _white(240, 80)
        cv2.rectangle(img, (0, 0), (239, 79), BLUE_BGR, -1)
        kind, _fill, _line, _radius, _line_px = classify_filled_shape(img)
        self.assertEqual(kind, "rect")


def _write_layout_files(td: Path, img: np.ndarray,
                        bbox: tuple, role: str = "internal"):
    src = td / "page_01.png"
    cv2.imwrite(str(src), img)
    inventory = [
        {"id": "t000", "type": "text", "text": "标题",
         "bbox": [10, 10, 60, 30], "confidence": 0.95},
        {"id": "v000", "type": "image", "bbox": list(bbox),
         "source": "cleaned", "role": role},
    ]
    inv_path = td / "inventory.json"
    inv_path.write_text(json.dumps(inventory), encoding="utf-8")
    return type("Args", (), {
        "inventory": str(inv_path),
        "source": str(src),
        "cleaned": str(src),
        "out_assets_dir": str(td / "assets" / "page_01"),
        "asset_prefix": "assets/page_01",
        "out_manifest": str(td / "m.json"),
        "out_layout": str(td / "l.json"),
        "slide_width_in": None,
        "slide_height_in": 7.5,
    })()


class BuilderIntegrationTests(unittest.TestCase):
    def test_rotated_rect_flows_to_layout_and_pptx(self) -> None:
        from layout.builder import LayoutBuilder
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        img = _white(320, 260)
        th = math.radians(30)
        c, s = math.cos(th), math.sin(th)
        corners = np.array(
            [(-100, -35), (100, -35), (100, 35), (-100, 35)], np.float64)
        R = np.array([[c, -s], [s, c]])
        pts = (corners @ R.T) + np.array([160, 130])
        cv2.fillPoly(img, [np.round(pts).astype(np.int32)], BLUE_BGR)
        # tight rotated bbox: 200x70 at 30° → 208x162 centred (160,130)
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            args = _write_layout_files(td, img, (56, 49, 264, 211))
            builder = LayoutBuilder(args)
            builder.build()
            builder.write()
            self.assertEqual(len(builder.front_shape_elements), 1)
            shape = builder.front_shape_elements[0]
            self.assertEqual(shape["shape"], "rect")
            self.assertAlmostEqual(shape["rotation"], 30.0, delta=3.0)
            layout = json.loads(
                (td / "l.json").read_text(encoding="utf-8"))
            rot = [el for el in layout["elements"]
                   if el.get("type") == "shape"]
            self.assertEqual(len(rot), 1)
            prs = Presentation(str(self._build_pptx(td)))
            autos = [sh for sh in prs.slides[0].shapes
                     if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
            self.assertEqual(len(autos), 1)
            self.assertAlmostEqual(abs(autos[0].rotation), 30.0, delta=3.0)

    @staticmethod
    def _build_pptx(td: Path) -> Path:
        from deck.build_pptx_from_layout import run as build_pptx
        out = td / "out.pptx"
        build_pptx(layout=str(td / "l.json"), out=str(out),
                   assets_root=str(td))
        return out

    def test_wide_ellipse_outline_becomes_oval_shape(self) -> None:
        from layout.builder import LayoutBuilder
        img = _white(500, 300)
        cv2.ellipse(img, (250, 150), (200, 90), 0, 0, 360,
                    (180, 100, 40), 4)
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            args = _write_layout_files(td, img, (50, 60, 450, 240),
                                       role="outline")
            builder = LayoutBuilder(args)
            builder.build()
            builder.write()
            self.assertEqual(len(builder.shape_elements), 1)
            shape = builder.shape_elements[0]
            self.assertEqual(shape["shape"], "oval")
            self.assertEqual(shape["box"], [50, 60, 400, 180])
            self.assertEqual(builder.image_elements, [])

    def test_parallelogram_element_has_no_radius_key(self) -> None:
        """Parallelogram/hexagon/pentagon elements must NOT carry
        radius: zeroing the MSO slant adjustment would deform them."""
        from layout.builder import LayoutBuilder
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
        img = _white(240, 140)
        pts = np.array([[50, 2], [230, 2], [200, 138], [20, 138]],
                       np.int32)
        cv2.fillPoly(img, [pts], GRAY_BGR)
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            args = _write_layout_files(td, img, (20, 2, 230, 138))
            builder = LayoutBuilder(args)
            builder.build()
            builder.write()
            self.assertEqual(len(builder.front_shape_elements), 1)
            shape = builder.front_shape_elements[0]
            self.assertEqual(shape["shape"], "parallelogram")
            self.assertNotIn("radius", shape)
            prs = Presentation(str(self._build_pptx(td)))
            autos = [sh for sh in prs.slides[0].shapes
                     if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
            self.assertEqual(len(autos), 1)
            self.assertEqual(autos[0].auto_shape_type,
                             MSO_SHAPE.PARALLELOGRAM)


if __name__ == "__main__":
    unittest.main()
