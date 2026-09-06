"""Native shape lift phase 2: block arrows, elbows, banners, dashed frames.

Second optimization pass over native shape recognition ("原生形状识别
优化"):

* classify_filled_shape lifts solid block arrows (right/left/up/down)
  with measured shaft/head proportions (arrow_geometry).
* classify_elbow_line lifts L-shaped connectors to 3-point polylines;
  the PPTX builder renders them as one native freeform.
* _polygon_candidate lifts right-pointing homeplate / chevron process
  banners and short-bottom trapezoids (trapezoid_down).
* _quad_ring_candidate / classify_outline_ring lift dashed frames to
  dashed_rect / dashed_round_rect (native dash styling).
* classify_connector_line lifts both-end arrows (arrow="both").
"""
from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

import cv2
import numpy as np

ROOT = Path(__file__).resolve().parents[1]
PAGE_DIR = ROOT / "scripts" / "page"
SCRIPTS_DIR = ROOT / "scripts"
for _p in (str(PAGE_DIR), str(SCRIPTS_DIR)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from layout.outline import (  # noqa: E402
    arrow_geometry,
    classify_connector_line,
    classify_elbow_line,
    classify_filled_shape,
    classify_outline_ring,
)

BLUE_BGR = (200, 120, 40)      # → #2878C8
GRAY_BGR = (90, 90, 90)        # → #5A5A5A


def _img(w: int, h: int) -> np.ndarray:
    return np.full((h, w, 3), 255, np.uint8)


def _right_arrow(w: int = 200, h: int = 100, shaft: int = 36,
                 head: int = 90, color=BLUE_BGR) -> np.ndarray:
    img = _img(w, h)
    cy = h // 2
    cv2.rectangle(img, (0, cy - shaft // 2), (w - head, cy + shaft // 2),
                  color, -1)
    cv2.fillPoly(img, [np.array(
        [[w - head, 0], [w, cy], [w - head, h]], np.int32)], color)
    return img


def _embed(mask: np.ndarray, color=BLUE_BGR) -> np.ndarray:
    h, w = mask.shape
    img = _img(w, h)
    img[mask] = color
    return img


class BlockArrowTests(unittest.TestCase):
    def test_right_arrow_lifts_with_geometry(self) -> None:
        kind, *_ = classify_filled_shape(_right_arrow())
        self.assertEqual(kind, "right_arrow")
        geom = arrow_geometry(_right_arrow())
        self.assertEqual(geom[0], "right_arrow")
        adj1, adj2 = geom[1], geom[2]
        # shaft 36 px of 100 px height; head 90 px of min side 100
        self.assertAlmostEqual(adj1, 0.36, delta=0.05)
        self.assertAlmostEqual(adj2, 0.90, delta=0.05)

    def test_all_four_orientations(self) -> None:
        ra = _right_arrow()
        self.assertEqual(classify_filled_shape(ra)[0], "right_arrow")
        self.assertEqual(classify_filled_shape(ra[:, ::-1])[0], "left_arrow")
        # spatial-only transpose: transpose maps the +x head to +y (down)
        self.assertEqual(classify_filled_shape(ra.transpose(1, 0, 2))[0],
                         "down_arrow")
        self.assertEqual(
            classify_filled_shape(ra.transpose(1, 0, 2)[::-1, :])[0],
            "up_arrow")

    def test_flag_banner_not_an_arrow(self) -> None:
        # full-height block on the right: no taper → not an arrow
        img = _img(200, 100)
        cv2.rectangle(img, (0, 30), (120, 70), BLUE_BGR, -1)
        cv2.rectangle(img, (120, 0), (199, 99), BLUE_BGR, -1)
        result = classify_filled_shape(img)
        self.assertTrue(result is None or result[0] != "right_arrow")

    def test_circle_square_capsule_not_arrows(self) -> None:
        circle = _img(160, 160)
        cv2.circle(circle, (80, 80), 78, BLUE_BGR, -1)
        self.assertEqual(classify_filled_shape(circle)[0], "oval")
        square = _img(100, 100)
        cv2.rectangle(square, (0, 0), (99, 99), BLUE_BGR, -1)
        self.assertEqual(classify_filled_shape(square)[0], "rect")
        cap = _img(300, 100)
        cv2.rectangle(cap, (50, 0), (250, 99), BLUE_BGR, -1)
        cv2.circle(cap, (50, 50), 50, BLUE_BGR, -1)
        cv2.circle(cap, (250, 50), 50, BLUE_BGR, -1)
        self.assertEqual(classify_filled_shape(cap)[0], "round_rect")


class BannerTests(unittest.TestCase):
    def test_homeplate_lifts(self) -> None:
        img = _img(240, 120)
        cv2.fillPoly(img, [np.array(
            [[0, 0], [170, 0], [239, 60], [170, 119], [0, 119]], np.int32)],
            BLUE_BGR)
        self.assertEqual(classify_filled_shape(img)[0], "homeplate")

    def test_chevron_lifts(self) -> None:
        img = _img(260, 120)
        cv2.fillPoly(img, [np.array(
            [[0, 0], [180, 0], [259, 60], [180, 119], [0, 119], [70, 60]],
            np.int32)], GRAY_BGR)
        self.assertEqual(classify_filled_shape(img)[0], "chevron")


class TrapezoidOrientationTests(unittest.TestCase):
    def test_short_bottom_trapezoid_lifts_down(self) -> None:
        img = _img(200, 120)
        cv2.fillPoly(img, [np.array(
            [[0, 0], [199, 0], [170, 119], [30, 119]], np.int32)], BLUE_BGR)
        self.assertEqual(classify_filled_shape(img)[0], "trapezoid_down")

    def test_short_top_trapezoid_unchanged(self) -> None:
        img = _img(200, 120)
        cv2.fillPoly(img, [np.array(
            [[30, 0], [170, 0], [199, 119], [0, 119]], np.int32)], BLUE_BGR)
        self.assertEqual(classify_filled_shape(img)[0], "trapezoid")


class DashedFrameTests(unittest.TestCase):
    @staticmethod
    def _dashed_frame(inset: int) -> np.ndarray:
        img = _img(300, 160)
        x1, y1, x2, y2 = inset, inset, 299 - inset, 159 - inset
        for x in range(x1 + 2, x2 - 1, 14):
            cv2.line(img, (x, y1), (min(x + 7, x2 - 1), y1), BLUE_BGR, 3)
            cv2.line(img, (x, y2), (min(x + 7, x2 - 1), y2), BLUE_BGR, 3)
        for y in range(y1 + 2, y2 - 1, 14):
            cv2.line(img, (x1, y), (x1, min(y + 7, y2 - 1)), BLUE_BGR, 3)
            cv2.line(img, (x2, y), (x2, min(y + 7, y2 - 1)), BLUE_BGR, 3)
        return img

    def test_dashed_frame_lifts_with_dash_kind(self) -> None:
        kind, fill, line, _r, thickness = classify_filled_shape(
            self._dashed_frame(0))
        self.assertEqual(kind, "dashed_rect")
        self.assertIsNone(fill)          # transparent interior
        self.assertEqual(line, "#2878C8")
        self.assertGreaterEqual(thickness, 1.5)

    def test_dashed_inset_frame_lifts(self) -> None:
        kind, *_ = classify_filled_shape(self._dashed_frame(6))
        self.assertEqual(kind, "dashed_rect")

    def test_dashed_outline_card_lifts(self) -> None:
        result = classify_outline_ring(self._dashed_frame(6),
                                       (0, 0, 300, 160))
        self.assertIsNotNone(result)
        self.assertEqual(result[0], "dashed_round_rect")

    def test_solid_frame_keeps_plain_kind(self) -> None:
        img = _img(300, 160)
        cv2.rectangle(img, (0, 0), (299, 159), BLUE_BGR, 5)
        self.assertEqual(classify_filled_shape(img)[0], "rect")
        self.assertEqual(classify_outline_ring(img, (0, 0, 300, 160))[0],
                         "round_rect")


class ElbowConnectorTests(unittest.TestCase):
    def _elbow(self, h_top: bool, v_left: bool) -> np.ndarray:
        img = _img(220, 160)
        y1, y2 = (0, 6) if h_top else (153, 159)
        x1, x2 = (0, 210) if v_left else (10, 219)
        hx1, hx2 = (0, 210) if v_left else (10, 219)
        vy1, vy2 = (0, 150) if h_top else (9, 159)
        vx1, vx2 = (0, 6) if v_left else (213, 219)
        cv2.rectangle(img, (hx1, y1), (hx2, y2), BLUE_BGR, -1)
        cv2.rectangle(img, (vx1, vy1), (vx2, vy2), BLUE_BGR, -1)
        return img

    def test_four_corner_orientations_lift(self) -> None:
        for h_top in (True, False):
            for v_left in (True, False):
                hit = classify_elbow_line(self._elbow(h_top, v_left))
                self.assertIsNotNone(hit, f"h_top={h_top} v_left={v_left}")
                pts, line, width, dash, arrow = hit
                self.assertEqual(len(pts), 6)
                self.assertEqual(line, "#2878C8")
                self.assertIsNone(dash)
                self.assertIsNone(arrow)
                self.assertGreaterEqual(width, 1.0)

    def test_z_shape_keeps_png(self) -> None:
        img = _img(220, 160)
        cv2.rectangle(img, (0, 0), (210, 6), BLUE_BGR, -1)
        cv2.rectangle(img, (204, 0), (210, 150), BLUE_BGR, -1)
        cv2.rectangle(img, (0, 150), (210, 156), BLUE_BGR, -1)
        self.assertIsNone(classify_elbow_line(img))

    def test_frame_and_straight_rejected(self) -> None:
        frame = _img(220, 160)
        cv2.rectangle(frame, (0, 0), (210, 150), BLUE_BGR, 5)
        self.assertIsNone(classify_elbow_line(frame))
        straight = _img(220, 160)
        cv2.line(straight, (10, 80), (210, 80), BLUE_BGR, 5)
        self.assertIsNone(classify_elbow_line(straight))


class BothEndArrowTests(unittest.TestCase):
    def test_double_arrowhead_lifts(self) -> None:
        img = _img(300, 40)
        cv2.rectangle(img, (0, 17), (299, 22), BLUE_BGR, -1)
        cv2.fillPoly(img, [np.array([[40, 8], [40, 31], [10, 20]],
                                    np.int32)], BLUE_BGR)
        cv2.fillPoly(img, [np.array([[260, 8], [260, 31], [290, 20]],
                                    np.int32)], BLUE_BGR)
        hit = classify_connector_line(img)
        self.assertIsNotNone(hit)
        self.assertEqual(hit[4], "both")


def _builder_args(td: Path, inventory: list) -> object:
    src = td / "page_01.png"
    return type("Args", (), {
        "inventory": str(td / "inventory.json"),
        "source": str(src),
        "cleaned": str(src),
        "out_assets_dir": str(td / "assets" / "page_01"),
        "asset_prefix": "assets/page_01",
        "out_manifest": str(td / "m.json"),
        "out_layout": str(td / "l.json"),
        "slide_width_in": None,
        "slide_height_in": 7.5,
    })()


class BuilderIntegrationTests(unittest.TestCase):
    def test_arrow_flows_to_layout_and_pptx_with_adjustments(self) -> None:
        from layout.builder import LayoutBuilder
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
        img = _img(320, 260)
        img[50:150, 50:250] = _right_arrow()
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            cv2.imwrite(str(td / "page_01.png"), img)
            (td / "inventory.json").write_text(json.dumps([
                {"id": "v000", "type": "image",
                 "bbox": [50, 50, 250, 150],
                 "source": "cleaned", "role": "internal"},
            ]), encoding="utf-8")
            builder = LayoutBuilder(_builder_args(td, []))
            builder.build()
            builder.write()
            self.assertEqual(len(builder.front_shape_elements), 1)
            el = builder.front_shape_elements[0]
            self.assertEqual(el["shape"], "right_arrow")
            self.assertEqual(len(el["adjustments"]), 2)
            self.assertEqual(builder.image_elements, [])
            from deck.build_pptx_from_layout import run as build_pptx
            out = td / "o.pptx"
            build_pptx(layout=str(td / "l.json"), out=str(out),
                       assets_root=str(td))
            prs = Presentation(str(out))
            autos = [sh for sh in prs.slides[0].shapes
                     if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
            self.assertEqual(len(autos), 1)
            self.assertEqual(autos[0].auto_shape_type, MSO_SHAPE.RIGHT_ARROW)
            self.assertAlmostEqual(float(autos[0].adjustments[0]),
                                   el["adjustments"][0], delta=0.01)

    def test_elbow_flows_to_freeform_line(self) -> None:
        from layout.builder import LayoutBuilder
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        img = _img(600, 400)
        cv2.rectangle(img, (300, 200), (560, 206), BLUE_BGR, -1)
        cv2.rectangle(img, (300, 200), (306, 360), BLUE_BGR, -1)
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            cv2.imwrite(str(td / "page_01.png"), img)
            (td / "inventory.json").write_text(json.dumps([
                {"id": "v000", "type": "image",
                 "bbox": [300, 200, 561, 361],
                 "source": "cleaned", "role": "connector"},
            ]), encoding="utf-8")
            builder = LayoutBuilder(_builder_args(td, []))
            builder.build()
            builder.write()
            lines = [el for el in builder.front_shape_elements
                     if el.get("type") == "line"]
            self.assertEqual(len(lines), 1)
            self.assertEqual(len(lines[0]["points"]), 6)
            from deck.build_pptx_from_layout import run as build_pptx
            out = td / "o.pptx"
            build_pptx(layout=str(td / "l.json"), out=str(out),
                       assets_root=str(td))
            prs = Presentation(str(out))
            frees = [sh for sh in prs.slides[0].shapes
                     if sh.shape_type == MSO_SHAPE_TYPE.FREEFORM]
            self.assertEqual(len(frees), 1)
            self.assertEqual(str(frees[0].line.color.rgb), "2878C8")

    def test_trapezoid_down_and_dashed_render_in_pptx(self) -> None:
        from deck.build_pptx_from_layout import Builder
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        layout = {
            "slide_size": {"width_in": 13.333, "height_in": 7.5},
            "source_width": 600, "source_height": 400,
            "background": "#FFFFFF",
            "elements": [
                {"type": "shape", "name": "td",
                 "shape": "trapezoid_down", "box": [10, 10, 100, 60],
                 "fill": "#2878C8", "line": "#2878C8", "line_width": 1},
                {"type": "shape", "name": "dd",
                 "shape": "dashed_round_rect", "box": [150, 10, 200, 100],
                 "fill": "transparent", "line": "#C8C8C8",
                 "line_width": 1, "radius": 0.12},
            ],
        }
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            (td / "l.json").write_text(json.dumps(layout), encoding="utf-8")
            b = Builder(layout, td / "o.pptx", td)
            b.build()
            prs = Presentation(str(td / "o.pptx"))
            shapes = list(prs.slides[0].shapes)
            self.assertEqual(shapes[0].auto_shape_type, MSO_SHAPE.TRAPEZOID)
            self.assertEqual(shapes[0].rotation, 180.0)
            self.assertEqual(shapes[1].auto_shape_type,
                             MSO_SHAPE.ROUNDED_RECTANGLE)
            self.assertEqual(shapes[1].line.dash_style,
                             MSO_LINE_DASH_STYLE.DASH)


if __name__ == "__main__":
    unittest.main()
