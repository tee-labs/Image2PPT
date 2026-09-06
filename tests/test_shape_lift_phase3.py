"""Native shape lift phase 3: stars, sharp frames, dashed arrows, strokes.

Third optimization pass over native shape recognition ("原生形状识别
优化"):

* classify_filled_shape lifts concave 4/5-point stars (sparkles, rating
  stars) via an alternating-radius structural test — MSO orientation
  gates keep only point-up lifts.
* classify_outline_ring lifts sharp-cornered frames as plain rect /
  dashed_rect (the old 0.08 minimum radius rounded corners the source
  does not have) and reroutes 50%-duty dash frames through the
  dash-bridging branch.
* classify_connector_line lifts dashed arrows (dash and arrow styling
  are independent layout keys); straightness is judged on the middle
  half of the axis so end arrowheads do not trip the perp-spread gate,
  and blunt end bulges (L-bend arms) still reject.
* Solid lifts carry a measured border-stroke width (line_px) instead of
  the builder's fixed 0.9 pt default.
"""
from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

import cv2
import numpy as np

ROOT = Path(__file__).resolve().parents[1]
PAGE_DIR = ROOT / "scripts" / "page"
SCRIPTS_DIR = ROOT / "scripts"
for _p in (str(PAGE_DIR), str(SCRIPTS_DIR)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from layout.outline import (  # noqa: E402
    classify_connector_line,
    classify_filled_shape,
    classify_outline_ring,
)

BLUE_BGR = (200, 120, 40)      # → #2878C8
GRAY_BGR = (60, 60, 60)        # → #3C3C3C


def _img(w: int, h: int) -> np.ndarray:
    return np.full((h, w, 3), 255, np.uint8)


def _star(cx: float, cy: float, r_out: float, r_in: float,
          points: int, rot_deg: float = -90.0) -> np.ndarray:
    pts = []
    for i in range(2 * points):
        r = r_out if i % 2 == 0 else r_in
        ang = np.deg2rad(rot_deg + i * 180.0 / points)
        pts.append([cx + r * np.cos(ang), cy + r * np.sin(ang)])
    return np.array(pts, np.int32)


class StarLiftTests(unittest.TestCase):
    def test_five_point_star_lifts(self) -> None:
        img = _img(120, 120)
        cv2.fillPoly(img, [_star(60, 62, 54, 21, 5)], GRAY_BGR)
        kind, fill, line, radius, line_px = classify_filled_shape(img)
        self.assertEqual(kind, "star_5")
        self.assertEqual(fill, "#3C3C3C")
        self.assertEqual(radius, 0.0)

    def test_four_point_sparkle_lifts(self) -> None:
        img = _img(120, 120)
        cv2.fillPoly(img, [_star(60, 60, 54, 13, 4)], BLUE_BGR)
        self.assertEqual(classify_filled_shape(img)[0], "star_4")

    def test_down_pointing_star_stays_png(self) -> None:
        # MSO star presets are point-up; a 180° flip would render at
        # the wrong angle, so the orientation gate keeps the PNG.
        img = _img(120, 120)
        cv2.fillPoly(img, [_star(60, 58, 54, 21, 5, rot_deg=90.0)],
                     GRAY_BGR)
        self.assertIsNone(classify_filled_shape(img))

    def test_side_pointing_star_stays_png(self) -> None:
        img = _img(120, 120)
        cv2.fillPoly(img, [_star(60, 60, 54, 21, 5, rot_deg=-180.0)],
                     GRAY_BGR)
        self.assertIsNone(classify_filled_shape(img))

    def test_gear_like_blob_stays_png(self) -> None:
        # Twelve alternating radii (6-point star / gear): unsupported
        # star kinds keep the pixel-perfect PNG path.
        img = _img(140, 140)
        cv2.fillPoly(img, [_star(70, 70, 62, 40, 6)], BLUE_BGR)
        self.assertIsNone(classify_filled_shape(img))

    def test_sunburst_inconsistent_arms_stay_png(self) -> None:
        # Alternating radii with wildly different arm lengths are not
        # stars.
        img = _img(160, 160)
        pts = []
        for i in range(10):
            r = (70, 60, 40, 55, 30)[i % 5] if i % 2 == 0 else 18
            ang = np.deg2rad(-90 + i * 36.0)
            pts.append([80 + r * np.cos(ang), 80 + r * np.sin(ang)])
        cv2.fillPoly(img, [np.array(pts, np.int32)], GRAY_BGR)
        self.assertIsNone(classify_filled_shape(img))

    def test_star_flows_to_layout_and_pptx(self) -> None:
        from layout.builder import LayoutBuilder
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
        img = _img(320, 260)
        cv2.fillPoly(img, [_star(160, 132, 95, 38, 5)], BLUE_BGR)
        with tempfile.TemporaryDirectory() as td_s:
            td = Path(td_s)
            cv2.imwrite(str(td / "page_01.png"), img)
            (td / "inventory.json").write_text(json.dumps([
                {"id": "v000", "type": "image",
                 "bbox": [60, 32, 260, 232],
                 "source": "cleaned", "role": "internal"},
            ]), encoding="utf-8")
            args = type("Args", (), {
                "inventory": str(td / "inventory.json"),
                "source": str(td / "page_01.png"),
                "cleaned": str(td / "page_01.png"),
                "out_assets_dir": str(td / "assets" / "page_01"),
                "asset_prefix": "assets/page_01",
                "out_manifest": str(td / "m.json"),
                "out_layout": str(td / "l.json"),
                "slide_width_in": None,
                "slide_height_in": 7.5,
            })()
            builder = LayoutBuilder(args)
            builder.build()
            builder.write()
            self.assertEqual(len(builder.front_shape_elements), 1)
            el = builder.front_shape_elements[0]
            self.assertEqual(el["shape"], "star_5")
            from deck.build_pptx_from_layout import run as build_pptx
            out = td / "o.pptx"
            build_pptx(layout=str(td / "l.json"), out=str(out),
                       assets_root=str(td))
            prs = Presentation(str(out))
            autos = [sh for sh in prs.slides[0].shapes
                     if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
            self.assertEqual(len(autos), 1)
            self.assertEqual(autos[0].auto_shape_type,
                             MSO_SHAPE.STAR_5_POINT)


class SharpFrameTests(unittest.TestCase):
    def test_sharp_outline_ring_is_rect(self) -> None:
        source = _img(400, 200)
        cv2.rectangle(source, (3, 3), (396, 196), GRAY_BGR, 3)
        self.assertEqual(classify_outline_ring(source, (0, 0, 400, 200)),
                         ("rect", 0.0))

    def test_sharp_dashed_outline_ring_is_dashed_rect(self) -> None:
        source = _img(300, 160)
        for x in range(2, 297, 14):
            cv2.line(source, (x, 2), (min(x + 6, 296), 2), GRAY_BGR, 3)
            cv2.line(source, (x, 157), (min(x + 6, 296), 157), GRAY_BGR, 3)
        for y in range(2, 157, 14):
            cv2.line(source, (2, y), (2, min(y + 6, 156)), GRAY_BGR, 3)
            cv2.line(source, (297, y), (297, min(y + 6, 156)), GRAY_BGR, 3)
        result = classify_outline_ring(source, (0, 0, 300, 160))
        self.assertIsNotNone(result)
        self.assertEqual(result[0], "dashed_rect")
        self.assertEqual(result[1], 0.0)

    def test_rounded_outline_ring_keeps_radius(self) -> None:
        # A visibly rounded frame (corner stand-off ≫ AA fudge) keeps
        # the round_rect lift with its measured radius.
        source = _img(300, 200)
        cv2.rectangle(source, (4, 4), (295, 195), GRAY_BGR, 3)
        cv2.rectangle(source, (16, 4), (283, 195), GRAY_BGR, -1)
        cv2.rectangle(source, (16, 4), (283, 195), (255, 255, 255), -1)
        source[:20, :20] = 255  # clear the filled corner
        cv2.circle(source, (16, 16), 12, (255, 255, 255), -1)
        result = classify_outline_ring(source, (0, 0, 300, 200))
        self.assertIsNotNone(result)
        self.assertEqual(result[0], "round_rect")
        self.assertGreaterEqual(result[1], 0.08)


class DashedArrowTests(unittest.TestCase):
    def _dashed_arrow(self, dashed: bool) -> np.ndarray:
        img = _img(220, 120)
        if dashed:
            for x in range(4, 190, 14):
                cv2.line(img, (x, 58), (x + 8, 58), GRAY_BGR, 2)
        else:
            cv2.line(img, (4, 58), (188, 58), GRAY_BGR, 2)
        cv2.fillPoly(img, [np.array(
            [[188, 46], [188, 70], [204, 70], [204, 80],
             [218, 58], [204, 36], [204, 46]], np.int32)], GRAY_BGR)
        return img

    def test_dashed_arrow_lifts_with_both_styles(self) -> None:
        hit = classify_connector_line(self._dashed_arrow(dashed=True))
        self.assertIsNotNone(hit)
        _pts, line, width, dash, arrow = hit
        self.assertEqual(dash, "dash")
        self.assertIn(arrow, ("start", "end"))
        self.assertEqual(line, "#3C3C3C")
        self.assertGreaterEqual(width, 1.0)

    def test_solid_arrow_has_no_dash(self) -> None:
        hit = classify_connector_line(self._dashed_arrow(dashed=False))
        self.assertIsNotNone(hit)
        self.assertIsNone(hit[3])
        self.assertIn(hit[4], ("start", "end"))

    def test_blunt_arm_bend_rejected(self) -> None:
        # Long-arm L: the bend sits outside the mid-band, so straightness
        # passes, but the blunt vertical arm must reject the straight
        # emit (the elbow classifier owns the stroke).
        img = _img(300, 140)
        cv2.line(img, (4, 20), (250, 20), GRAY_BGR, 3)
        cv2.line(img, (250, 20), (250, 130), GRAY_BGR, 3)
        self.assertIsNone(classify_connector_line(img))


class BorderStrokeWidthTests(unittest.TestCase):
    def test_bordered_solid_measures_stroke(self) -> None:
        img = _img(160, 120)
        cv2.rectangle(img, (6, 6), (153, 113), BLUE_BGR, -1)
        for t in range(4):
            cv2.rectangle(img, (6 + t, 6 + t), (153 - t, 113 - t),
                          GRAY_BGR, 1)
        kind, fill, line, _radius, line_px = classify_filled_shape(img)
        self.assertEqual(kind, "rect")
        self.assertEqual(fill, "#2878C8")
        self.assertEqual(line, "#3C3C3C")
        self.assertGreaterEqual(line_px, 2.0)
        self.assertLessEqual(line_px, 5.0)

    def test_borderless_solid_keeps_zero(self) -> None:
        img = _img(160, 120)
        cv2.rectangle(img, (6, 6), (153, 113), BLUE_BGR, -1)
        _kind, _fill, _line, _radius, line_px = classify_filled_shape(img)
        self.assertEqual(line_px, 0.0)


if __name__ == "__main__":
    unittest.main()
