"""
Build a minimal "image-as-background + editable text" PPTX from an existing
run's erased page image and OCR JSON.

This is intentionally simpler than build_deck.py: no icon/shape/table
reconstruction, just the text-erased page picture plus an editable text box
per OCR entry. Text size and color are estimated from the source pixels.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Emu, Inches, Pt


def estimate_text_color(arr: np.ndarray, x1: int, y1: int, x2: int, y2: int) -> RGBColor:
    h_img, w_img = arr.shape[:2]
    x1 = max(0, min(x1, w_img - 1))
    x2 = max(x1 + 1, min(x2, w_img))
    y1 = max(0, min(y1, h_img - 1))
    y2 = max(y1 + 1, min(y2, h_img))
    crop = arr[y1:y2, x1:x2]
    if crop.size == 0:
        return RGBColor(0, 0, 0)
    h, w = crop.shape[:2]
    border = np.concatenate(
        [crop[0, :, :], crop[-1, :, :], crop[:, 0, :], crop[:, -1, :]]
    )
    bg = np.median(border, axis=0)
    if bg.mean() < 120:
        return RGBColor(255, 255, 255)
    gray = crop.mean(axis=2)
    thresh = np.percentile(gray, 8)
    mask = gray <= thresh
    pixels = crop[mask]
    if len(pixels) == 0:
        return RGBColor(0, 0, 0)
    r, g, b = (int(c) for c in np.median(pixels, axis=0))
    return RGBColor(r, g, b)


def build(
    work_dir: Path,
    out_path: Path,
    slide_width_in: float = 13.333,
    font_scale: float = 0.68,
) -> Path:
    erased = work_dir / "inventory" / "page_01.clean.text_only.png"
    if not erased.exists():
        erased = work_dir / "inventory" / "page_01.clean.png"
    original = work_dir / "source" / "page_01.png"
    ocr_path = work_dir / "ocr" / "page_01.ocr.json"

    if not erased.exists():
        raise FileNotFoundError(f"missing erased background: {erased}")
    if not original.exists():
        raise FileNotFoundError(f"missing source image: {original}")
    if not ocr_path.exists():
        raise FileNotFoundError(f"missing OCR JSON: {ocr_path}")

    with Image.open(original) as im:
        W, H = im.size
        orig_arr = np.array(im.convert("RGB"))

    ocr = json.loads(ocr_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_width_in * H / W)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide.shapes.add_picture(
        str(erased), 0, 0, width=prs.slide_width, height=prs.slide_height
    )

    emu_per_px = prs.slide_width / W

    for entry in ocr:
        text = (entry.get("text") or "").strip()
        if not text:
            continue
        x1, y1, x2, y2 = entry["x1"], entry["y1"], entry["x2"], entry["y2"]
        if x2 <= x1 or y2 <= y1:
            continue

        left = int(x1 * emu_per_px)
        top = int(y1 * emu_per_px)
        width = max(int((x2 - x1) * emu_per_px), Emu(91440))
        height = max(int((y2 - y1) * emu_per_px), Emu(91440))

        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        h_px = y2 - y1
        pt_size = max(8, int(round(h_px * font_scale)))
        run.font.size = Pt(pt_size)
        run.font.color.rgb = estimate_text_color(orig_arr, x1, y1, x2, y2)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument(
        "--work-dir",
        required=True,
        type=Path,
        help="Existing run directory produced by convert.py / build_deck.py.",
    )
    ap.add_argument(
        "--out",
        type=Path,
        default=None,
        help="Output PPTX path. Defaults to <work-dir>/slides_simple.pptx.",
    )
    ap.add_argument("--slide-width", type=float, default=13.333)
    ap.add_argument("--font-scale", type=float, default=0.68)
    args = ap.parse_args()

    out = args.out or (args.work_dir / "slides_simple.pptx")
    path = build(args.work_dir, out, args.slide_width, args.font_scale)
    print(path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
