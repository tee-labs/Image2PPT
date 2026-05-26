#!/usr/bin/env python
"""Build an editable PPTX from a slide layout JSON file."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

SCRIPTS_ROOT = Path(__file__).resolve().parents[1]
DECK_DIR = Path(__file__).resolve().parent
for _p in (DECK_DIR, SCRIPTS_ROOT):
    if str(_p) not in sys.path:
        sys.path.insert(0, str(_p))

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    raise SystemExit(
        "Missing python-pptx. Install it in the active Python environment with: "
        "python -m pip install python-pptx"
    ) from exc

from text_safety import ppt_safe_text  # noqa: E402

import text_finalizers  # noqa: E402
# Re-export so existing `from build_pptx_from_layout import apply_*` paths
# (used by tests + external scripts) keep working after the move.
from text_finalizers import (  # noqa: E402, F401
    apply_class_alignment,
    apply_size_unification,
    apply_title_centering,
)


def default_ppt_font() -> str:
    return "Microsoft YaHei"


DEFAULT_FONT = default_ppt_font()


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build editable PPTX from layout JSON.")
    parser.add_argument("--layout", required=True, help="Layout JSON path.")
    parser.add_argument("--out", required=True, help="Output .pptx path.")
    parser.add_argument("--assets-root", help="Base directory for relative image paths.")
    return parser.parse_args()


def rgb(value: Any, default: RGBColor | None = None) -> RGBColor | None:
    if value is None:
        return default
    if isinstance(value, str):
        text = value.strip()
        if text.lower() in {"none", "transparent", "background"}:
            return None
        text = text.lstrip("#")
        if len(text) == 6:
            return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        return RGBColor(int(value[0]), int(value[1]), int(value[2]))
    raise ValueError(f"Unsupported color: {value!r}")


def rgb_hex(value: Any) -> str | None:
    if value is None:
        return None
    if isinstance(value, str):
        text = value.strip()
        if text.lower() in {"none", "transparent", "background"}:
            return None
        text = text.lstrip("#")
        if len(text) == 6:
            return text.upper()
        return None
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        return f"{int(value[0]):02X}{int(value[1]):02X}{int(value[2]):02X}"
    return None


def align(value: str | None) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get((value or "center").lower(), PP_ALIGN.CENTER)


def valign(value: str | None) -> MSO_ANCHOR:
    return {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get((value or "middle").lower(), MSO_ANCHOR.MIDDLE)


def shape_type(value: str | None) -> MSO_SHAPE:
    return {
        "rect": MSO_SHAPE.RECTANGLE,
        "rectangle": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "ellipse": MSO_SHAPE.OVAL,
        "diamond": MSO_SHAPE.DIAMOND,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "trapezoid": MSO_SHAPE.TRAPEZOID,
    }.get((value or "rect").lower(), MSO_SHAPE.RECTANGLE)


def dash_style(value: str | None):
    if not value:
        return None
    return {
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "dashed": MSO_LINE_DASH_STYLE.DASH,
        "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dotted": MSO_LINE_DASH_STYLE.ROUND_DOT,
    }.get(value.lower())


class Builder:
    def __init__(self, layout: dict[str, Any], out: Path, assets_root: Path):
        self.layout = layout
        self.out = out
        self.assets_root = assets_root
        self.default_px_w = float(layout.get("source_width") or layout.get("canvas", {}).get("width") or 1182)
        self.default_px_h = float(layout.get("source_height") or layout.get("canvas", {}).get("height") or 665)
        slide_size = layout.get("slide_size", {})
        self.slide_w_in = float(slide_size.get("width_in", 13.333333))
        self.slide_h_in = float(slide_size.get("height_in", 7.5))
        self.scale_in_per_px = min(
            self.slide_w_in / self.default_px_w,
            self.slide_h_in / self.default_px_h,
        )
        self.offset_x_in = (self.slide_w_in - self.default_px_w * self.scale_in_per_px) / 2.0
        self.offset_y_in = (self.slide_h_in - self.default_px_h * self.scale_in_per_px) / 2.0
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_w_in)
        self.prs.slide_height = Inches(self.slide_h_in)

    def x(self, value: float):
        return Inches(self.offset_x_in + float(value) * self.scale_in_per_px)

    def y(self, value: float):
        return Inches(self.offset_y_in + float(value) * self.scale_in_per_px)

    def w(self, value: float):
        return Inches(float(value) * self.scale_in_per_px)

    def h(self, value: float):
        return Inches(float(value) * self.scale_in_per_px)

    def set_slide_coordinate_space(self, spec: dict[str, Any]) -> None:
        """Map this slide's source pixels into the deck canvas.

        PowerPoint uses one page size for the whole deck. When source
        screenshots have mixed aspect ratios, each slide is therefore
        letterboxed into the deck canvas instead of being stretched using
        the first slide's pixel dimensions.
        """
        px_w = float(spec.get("source_width") or self.default_px_w)
        px_h = float(spec.get("source_height") or self.default_px_h)
        if px_w <= 0 or px_h <= 0:
            px_w, px_h = self.default_px_w, self.default_px_h
        self.scale_in_per_px = min(
            self.slide_w_in / px_w,
            self.slide_h_in / px_h,
        )
        self.offset_x_in = (self.slide_w_in - px_w * self.scale_in_per_px) / 2.0
        self.offset_y_in = (self.slide_h_in - px_h * self.scale_in_per_px) / 2.0

    def set_font(self, run, spec: dict[str, Any]) -> None:
        font = spec.get("font") or spec.get("font_name") or DEFAULT_FONT
        run.font.name = font
        size = int(round(float(spec.get("size", spec.get("font_size", 18)))))
        run.font.size = Pt(max(1, size))
        run.font.bold = bool(spec.get("bold", False))
        run.font.italic = bool(spec.get("italic", False))
        col = rgb(spec.get("color", "#111111"))
        if col is not None:
            run.font.color.rgb = col
        r_pr = run._r.get_or_add_rPr()
        # Prevent PowerPoint from drawing red spell-check squiggles over
        # editable OCR text. These marks are UI proofing overlays, not part
        # of the original slide, and are especially noisy for Chinese text
        # mixed with symbols or English terms.
        r_pr.set("lang", str(spec.get("lang") or "zh-CN"))
        r_pr.set("noProof", "1")
        r_pr.set("dirty", "0")
        char_spacing = spec.get("char_spacing")
        if char_spacing is None:
            char_spacing = spec.get("character_spacing")
        if char_spacing is not None:
            r_pr.set("spc", str(int(round(float(char_spacing)))))
        for tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
            face = r_pr.find(qn(tag))
            if face is None:
                face = OxmlElement(tag)
                r_pr.append(face)
            face.set("typeface", font)

    def strip_list_marker_text(self, el: dict[str, Any], text: str) -> str:
        spec = el.get("list") or {}
        marker_chars = int(spec.get("marker_chars") or 0)
        if marker_chars <= 0:
            return text
        return text[marker_chars:].lstrip()

    def strip_list_marker_runs(self, el: dict[str, Any],
                               runs: list[dict[str, Any]]) -> list[dict[str, Any]]:
        spec = el.get("list") or {}
        remaining = int(spec.get("marker_chars") or 0)
        if remaining <= 0:
            return runs
        out: list[dict[str, Any]] = []
        stripped_leading_space = False
        for run in runs:
            text = ppt_safe_text(run.get("text", ""))
            if remaining:
                if len(text) <= remaining:
                    remaining -= len(text)
                    continue
                text = text[remaining:]
                remaining = 0
            if not stripped_leading_space:
                text = text.lstrip()
                stripped_leading_space = True
            if not text:
                continue
            new_run = dict(run)
            new_run["text"] = text
            out.append(new_run)
        return out

    def apply_paragraph_list_style(self, paragraph, el: dict[str, Any],
                                   box_left_px: float, box_width_px: float) -> None:
        spec = el.get("list") or {}
        if not spec:
            return
        p_pr = paragraph._p.get_or_add_pPr()
        for tag in ("a:buNone", "a:buFont", "a:buClr", "a:buChar", "a:buAutoNum"):
            child = p_pr.find(qn(tag))
            if child is not None:
                p_pr.remove(child)

        marker_x = float(spec.get("marker_x", box_left_px))
        body_x = float(spec.get("body_x", marker_x + 14.0))
        body_indent_px = body_x - float(box_left_px)
        marker_indent_px = marker_x - float(box_left_px)
        body_indent_px = max(6.0, min(float(box_width_px) - 1.0,
                                      body_indent_px))
        marker_indent_px = max(0.0, min(body_indent_px - 2.0,
                                        marker_indent_px))
        p_pr.set("marL", str(int(self.w(body_indent_px))))
        p_pr.set("indent", str(int(self.w(marker_indent_px - body_indent_px))))

        font = el.get("font") or el.get("font_name") or DEFAULT_FONT
        bu_font = OxmlElement("a:buFont")
        bu_font.set("typeface", font)
        p_pr.append(bu_font)

        hex_color = rgb_hex(el.get("color"))
        if hex_color:
            bu_clr = OxmlElement("a:buClr")
            srgb = OxmlElement("a:srgbClr")
            srgb.set("val", hex_color)
            bu_clr.append(srgb)
            p_pr.append(bu_clr)

        if spec.get("kind") == "ordered":
            bu_auto = OxmlElement("a:buAutoNum")
            bu_auto.set("type", str(spec.get("auto_type") or "arabicPeriod"))
            bu_auto.set("startAt", str(int(spec.get("start") or 1)))
            p_pr.append(bu_auto)
        else:
            bu_char = OxmlElement("a:buChar")
            bu_char.set("char", str(spec.get("marker") or "•"))
            p_pr.append(bu_char)

    def image_path(self, value: str) -> Path:
        path = Path(value)
        if not path.is_absolute():
            path = self.assets_root / path
        return path

    def add_text(self, slide, el: dict[str, Any]) -> None:
        left, top, width, height = el["box"]
        shape = slide.shapes.add_textbox(self.x(left), self.y(top), self.w(width), self.h(height))
        shape.name = el.get("name", "text")
        # Ensure transparent fill and no border so text doesn't cover icons/images below
        shape.fill.background()
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = False
        margin = float(el.get("margin", 0))
        tf.margin_left = self.w(margin)
        tf.margin_right = self.w(margin)
        tf.margin_top = self.h(margin)
        tf.margin_bottom = self.h(margin)
        tf.vertical_anchor = valign(el.get("valign"))
        # In-bbox color variation (e.g. a red statistic embedded in a
        # gray sentence) arrives as `runs` — a list of {text, color}
        # dicts. Each run becomes a separate python-pptx run within the
        # same paragraph, with the element-level font/size/bold but its
        # own color. When `runs` is absent, the legacy single-run path
        # is used; size and bold are always inherited from `el`.
        runs_spec = el.get("runs")
        if runs_spec:
            runs_spec = self.strip_list_marker_runs(el, runs_spec)
            # Split runs across explicit \n line breaks: each line gets
            # its own paragraph, runs within a line stay inline.
            paragraph_runs: list[list[dict]] = [[]]
            for r in runs_spec:
                segments = ppt_safe_text(r.get("text", "")).split("\n")
                for s_idx, seg in enumerate(segments):
                    if s_idx > 0:
                        paragraph_runs.append([])
                    if seg:
                        paragraph_runs[-1].append({
                            "text": seg,
                            "color": r.get("color"),
                            "size": r.get("size"),
                            "bold": r.get("bold"),
                        })
            for p_idx, line_runs in enumerate(paragraph_runs):
                p = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
                p.alignment = align(el.get("align"))
                p.line_spacing = float(el.get("line_spacing", 1.05))
                if p_idx == 0:
                    self.apply_paragraph_list_style(p, el, left, width)
                if not line_runs:
                    # Empty paragraph (consecutive newlines); add a blank
                    # run with the element default color so the line keeps
                    # its height.
                    run = p.add_run()
                    run.text = ""
                    self.set_font(run, el)
                    continue
                for r in line_runs:
                    run = p.add_run()
                    run.text = r["text"]
                    run_spec = dict(el)  # inherit size/bold/font
                    if r.get("color"):
                        run_spec["color"] = r["color"]
                    if r.get("size"):
                        run_spec["size"] = r["size"]
                    if r.get("bold") is not None:
                        run_spec["bold"] = bool(r["bold"])
                    self.set_font(run, run_spec)
        else:
            text = self.strip_list_marker_text(
                el, ppt_safe_text(el.get("text", "")))
            lines = text.split("\n")
            for idx, line in enumerate(lines):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.alignment = align(el.get("align"))
                p.line_spacing = float(el.get("line_spacing", 1.05))
                if idx == 0:
                    self.apply_paragraph_list_style(p, el, left, width)
                run = p.add_run()
                run.text = line
                self.set_font(run, el)

    def add_image(self, slide, el: dict[str, Any]) -> None:
        left, top, width, height = el["box"]
        pic = slide.shapes.add_picture(
            str(self.image_path(el["path"])),
            self.x(left),
            self.y(top),
            width=self.w(width),
            height=self.h(height),
        )
        pic.name = el.get("name", Path(el["path"]).stem)
        if "rotation" in el:
            pic.rotation = float(el["rotation"])

    def add_shape(self, slide, el: dict[str, Any]) -> None:
        left, top, width, height = el["box"]
        shape = slide.shapes.add_shape(
            shape_type(el.get("shape")),
            self.x(left),
            self.y(top),
            self.w(width),
            self.h(height),
        )
        shape.name = el.get("name", el.get("shape", "shape"))
        fill = rgb(el.get("fill"))
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        line = rgb(el.get("line"), RGBColor(0, 0, 0))
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
            shape.line.width = Pt(float(el.get("line_width", 1)))
        if len(shape.adjustments) and "radius" in el:
            shape.adjustments[0] = float(el["radius"])
        if "rotation" in el:
            shape.rotation = float(el["rotation"])
        if "shadow" in el:
            shape.shadow.inherit = bool(el["shadow"])
        elif el.get("shadow_off", True):
            shape.shadow.inherit = False

    def add_table(self, slide, el: dict[str, Any]) -> None:
        """Render a native PowerPoint table.

        Layout schema:
            {
              "type": "table",
              "box": [x, y, w, h],            // in source-image px
              "rows": 4, "cols": 3,           // grid size
              "cells": [
                  {"row": 0, "col": 0,
                   "rowspan": 1, "colspan": 1,
                   "text": "...",
                   "bold": false,
                   "color": "#222222",
                   "fill": "#FFFFFF",         // optional cell fill
                   "align": "left"|"center"|"right",
                   "valign": "top"|"middle"|"bottom"
                  }, ...]
              "font": "Microsoft YaHei",      // default for all cells
              "size": 14                      // default for all cells
            }

        Cell entries are sparse — any (row, col) not listed renders as an
        empty cell with default style. Merged cells (rowspan/colspan>1)
        use python-pptx's `merge` API and stash the text on the top-left
        cell, mirroring SLANet's structure output.
        """
        left, top, width, height = el["box"]
        rows = int(el.get("rows", 0))
        cols = int(el.get("cols", 0))
        if rows < 1 or cols < 1:
            return
        shape = slide.shapes.add_table(
            rows, cols, self.x(left), self.y(top),
            self.w(width), self.h(height),
        )
        shape.name = el.get("name", "table")
        tbl = shape.table
        for i, col_w in enumerate(el.get("col_widths", [])[:cols]):
            tbl.columns[i].width = self.w(float(col_w))
        for i, row_h in enumerate(el.get("row_heights", [])[:rows]):
            tbl.rows[i].height = self.h(float(row_h))

        default_font = el.get("font", DEFAULT_FONT)
        default_size = float(el.get("size", 14))
        default_color = el.get("color", "#222222")
        default_align = el.get("align", "left")
        default_valign = el.get("valign", "middle")

        # First apply merges so subsequent cell writes target the right
        # logical cell. python-pptx merges by calling `cell.merge(other)`
        # — we merge from top-left corner of each span outward.
        for c in el.get("cells", []):
            rs = int(c.get("rowspan", 1))
            cs = int(c.get("colspan", 1))
            if rs <= 1 and cs <= 1:
                continue
            r0, c0 = int(c["row"]), int(c["col"])
            r1, c1 = r0 + rs - 1, c0 + cs - 1
            if r1 >= rows or c1 >= cols:
                continue
            try:
                tbl.cell(r0, c0).merge(tbl.cell(r1, c1))
            except (ValueError, KeyError):
                # python-pptx raises when a sub-cell is already part of
                # another merge. Skip overlapping spans rather than
                # aborting the slide.
                continue

        for c in el.get("cells", []):
            r0, c0 = int(c["row"]), int(c["col"])
            if r0 >= rows or c0 >= cols:
                continue
            cell = tbl.cell(r0, c0)
            text = str(c.get("text", "")).strip()
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            margin = float(c.get("margin", el.get("margin", 1.0)))
            tf.margin_left = self.w(margin)
            tf.margin_right = self.w(margin)
            tf.margin_top = self.h(margin)
            tf.margin_bottom = self.h(margin)
            p = tf.paragraphs[0]
            p.alignment = align(c.get("align", default_align))
            run = p.add_run()
            run.text = text
            self.set_font(run, {
                "font": c.get("font", default_font),
                "size": c.get("size", default_size),
                "bold": c.get("bold", False),
                "italic": c.get("italic", False),
                "color": c.get("color", default_color),
            })
            tf.vertical_anchor = valign(c.get("valign", default_valign))
            fill = rgb(c.get("fill"))
            if fill is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill

    def add_line(self, slide, el: dict[str, Any]) -> None:
        if "points" in el:
            x1, y1, x2, y2 = el["points"]
        else:
            left, top, width, height = el["box"]
            x1, y1, x2, y2 = left, top, left + width, top + height
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, self.x(x1), self.y(y1), self.x(x2), self.y(y2))
        line.name = el.get("name", "line")
        col = rgb(el.get("line", el.get("color", "#000000")))
        if col is not None:
            line.line.color.rgb = col
        line.line.width = Pt(float(el.get("line_width", el.get("weight", 1))))
        dash = dash_style(el.get("dash"))
        if dash:
            line.line.dash_style = dash

    def add_slide(self, spec: dict[str, Any]) -> None:
        self.set_slide_coordinate_space(spec)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = rgb(spec.get("background", self.layout.get("background", "#FFFFFF")))
        if bg is not None:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg
        elements = spec.get("elements", [])
        text_finalizers.apply_all(elements)
        for el in elements:
            kind = el.get("type", "shape").lower()
            if kind == "text":
                self.add_text(slide, el)
            elif kind == "image":
                self.add_image(slide, el)
            elif kind == "shape":
                self.add_shape(slide, el)
            elif kind == "line":
                self.add_line(slide, el)
            elif kind == "table":
                self.add_table(slide, el)
            else:
                raise ValueError(f"Unsupported element type: {kind}")

    def build(self) -> None:
        slides = self.layout.get("slides")
        if not slides:
            slides = [{"elements": self.layout.get("elements", [])}]
        for spec in slides:
            self.add_slide(spec)
        self.out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(self.out)


def run(*, layout: str, out: str,
        assets_root: str | None = None) -> None:
    """Programmatic entry — same contract as the CLI flags."""
    layout_path = Path(layout)
    layout_data = json.loads(layout_path.read_text(encoding="utf-8-sig"))
    root = Path(assets_root) if assets_root else layout_path.parent
    Builder(layout_data, Path(out), root).build()


def main() -> None:
    args = parse_args()
    run(layout=args.layout, out=args.out, assets_root=args.assets_root)
    print(args.out)


if __name__ == "__main__":
    main()
